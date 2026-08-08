import io
import math
from typing import Callable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from app.infographic_models import (
    BULLET_MAX_COUNT,
    COMPARISON_MAX_COLUMNS,
    COMPARISON_POINT_COUNT,
    HUB_SPOKE_ITEM_COUNT,
    MATRIX_ITEM_COUNT,
    MATRIX_QUADRANT_COUNT,
    PYRAMID_MAX_PILLARS,
    PYRAMID_MIN_PILLARS,
    ROADMAP_COLUMN_COUNT,
    ROADMAP_ITEM_COUNT,
    RACI_MAX_ROWS,
    TIMELINE_MAX_MILESTONES,
    TITLE_HIGHLIGHT_MAX,
    VALUE_PROP_MAX_ITEMS,
    AgendaSlide,
    BulletSummarySlide,
    FeatureStory,
    HubSpokeItem,
    InfographicComparison,
    InfographicDiagram,
    InfographicHubSpoke,
    InfographicMatrix,
    InfographicPyramid,
    InfographicRoadmap,
    InfographicTimeline,
    InfographicWheel,
    MatrixQuadrant,
    PositioningStatementSlide,
    PyramidPillar,
    RaciChartSlide,
    RoadmapColumn,
    TitleSlide,
    ValuePropositionSlide,
    WheelItem,
)

# Brand palette (matches frontend/src/lib/themes.ts's fidelity-green theme):
# primary green + gold accent lead every set, extended with validated
# green/gold-family hues (dataviz skill's six-check categorical validator,
# adjacent-pairs mode -- the correct test here since these are always
# consumed as a sequential wedge/list/column index, never freely paired,
# and every segment carries its own text label regardless).
BRAND_COLORS = ["00754A", "C9A227", "14B0A0", "8A5A0A", "4A9020", "5CC4A8"]
BRAND_COLORS_DARK = ["005F3C", "A4841F", "109083", "714908", "3C761A", "4BA089"]

WEDGE_COLORS = BRAND_COLORS[:5]
LIST_ACCENT_COLORS = BRAND_COLORS_DARK[:5]
COLUMN_COLORS = BRAND_COLORS[:4]

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

WHEEL_CX_IN = 3.15
WHEEL_CY_IN = 3.9
R_OUT_IN = 2.35
R_IN_IN = 0.95
LABEL_RADIUS_IN = (R_OUT_IN + R_IN_IN) / 2

LIST_X_IN = 7.75
LIST_W_IN = 4.9
LIST_H_IN = 1.1
LIST_TOP_IN = 0.55
LIST_GAP_IN = 1.38

COMPARISON_TITLE_Y_IN = 0.55
COMPARISON_COL_TOP_IN = 1.35
COMPARISON_MARGIN_IN = 0.75
COMPARISON_COL_GAP_IN = 0.35
COMPARISON_BAR_H_IN = 0.7

ROADMAP_TITLE_Y_IN = 0.5
ROADMAP_RAIL_Y_IN = 1.2
ROADMAP_COL_TOP_IN = 1.6
ROADMAP_MARGIN_IN = 0.9
ROADMAP_COL_GAP_IN = 0.4
ROADMAP_BAR_H_IN = 0.65
# Same hue (brand primary green), fading light -- closer horizons read
# visually "stronger." Validated as an ordinal ramp (dataviz skill:
# monotone lightness, light end still clears 2:1 contrast).
ROADMAP_COLORS = ["0A4A30", "00754A", "5CAE82"]
ROADMAP_TEXT_COLORS = ["FFFFFF", "FFFFFF", "1F2622"]

PYRAMID_HEADLINE_Y_IN = 1.0
PYRAMID_HEADLINE_H_IN = 1.5
PYRAMID_TOP_Y_IN = 2.35
PYRAMID_BOTTOM_Y_IN = 7.0
PYRAMID_CX_IN = SLIDE_W_IN / 2
PYRAMID_APEX_W_IN = 1.8
PYRAMID_BASE_W_IN = 8.5
PYRAMID_GAP_IN = 0.08

TIMELINE_TITLE_Y_IN = 0.5
TIMELINE_LINE_Y_IN = 3.9
TIMELINE_LEFT_IN = 1.1
TIMELINE_RIGHT_IN = 12.2
TIMELINE_CARD_W_IN = 1.95
TIMELINE_CARD_H_IN = 1.7
TIMELINE_CONNECTOR_H_IN = 0.35
TIMELINE_COLORS = BRAND_COLORS

BULLET_TITLE_Y_IN = 0.6
BULLET_LIST_TOP_IN = 1.7
BULLET_LIST_LEFT_IN = 1.5
BULLET_LIST_WIDTH_IN = 10.3
BULLET_ROW_H_IN = 0.85
BULLET_ACCENT_COLOR = "00754A"

MATRIX_TITLE_Y_IN = 0.5
MATRIX_GRID_TOP_IN = 1.3
MATRIX_GRID_LEFT_IN = 1.6
MATRIX_GRID_W_IN = 11.0
MATRIX_GRID_H_IN = 5.4
MATRIX_GRID_GAP_IN = 0.25
MATRIX_BAR_H_IN = 0.55

STORY_HEADLINE_Y_IN = 0.9
STORY_HEADLINE_H_IN = 1.2
STORY_PANEL_TOP_IN = 2.5
STORY_PANEL_H_IN = 4.3
STORY_MARGIN_IN = 0.8
STORY_GAP_IN = 0.5
STORY_BAR_H_IN = 0.6
# Muted bronze (friction) -> brand green (the product, building the fix) ->
# gold (the payoff), reinforcing the arc's tension even before the text is
# read, while staying inside the brand's green/gold family.
STORY_COLORS = {"problem": "8A5A0A", "solution": "00754A", "impact": "C9A227"}

HUB_SPOKE_CX_IN = SLIDE_W_IN / 2
HUB_SPOKE_CY_IN = 3.9
HUB_SPOKE_HUB_R_IN = 1.05
HUB_SPOKE_RING_R_IN_IN = 1.15
HUB_SPOKE_RING_R_OUT_IN = 1.4
HUB_SPOKE_RING_HALF_SPAN_DEG = 15
HUB_SPOKE_CARD_W_IN = 4.0
HUB_SPOKE_CARD_H_IN = 1.3
HUB_SPOKE_CARD_ROW_GAP_IN = 0.35
HUB_SPOKE_SPOKE_GAP_IN = 0.6
HUB_SPOKE_CAP_W_IN = 1.1
# 6 distinct hues -- the reference stock template this was adapted from
# reused one color across 2 of its 6 segments, which we deliberately fixed
# here, now drawn from the validated brand palette (see BRAND_COLORS above).
HUB_SPOKE_COLORS = BRAND_COLORS

TITLE_TOP_BAR_H_IN = 0.14
TITLE_HEADLINE_Y_IN = 3.0
TITLE_HEADLINE_H_IN = 1.1
TITLE_RULE_Y_IN = 3.82
TITLE_RULE_W_IN = 2.2
TITLE_RULE_H_IN = 0.035
TITLE_SUBTITLE_Y_IN = 4.35
TITLE_SUBTITLE_W_IN = 9.5
TITLE_SUBTITLE_H_IN = 1.0
TITLE_CHIP_Y_IN = 6.15
TITLE_CHIP_W_IN = 2.05
TITLE_CHIP_H_IN = 0.5
TITLE_CHIP_GAP_IN = 0.22

AGENDA_TITLE_Y_IN = 0.6
AGENDA_LIST_TOP_IN = 1.7
AGENDA_LIST_LEFT_IN = 1.5
AGENDA_LIST_WIDTH_IN = 10.3
AGENDA_ROW_H_IN = 0.66
AGENDA_CHIP_SIZE_IN = 0.5

VALUE_PROP_TITLE_Y_IN = 0.5
VALUE_PROP_PANEL_TOP_IN = 1.25
VALUE_PROP_PANEL_H_IN = 5.8
VALUE_PROP_PANEL_W_IN = 5.5
VALUE_PROP_LEFT_X_IN = 0.7
VALUE_PROP_RIGHT_X_IN = 7.15
VALUE_PROP_HEADER_H_IN = 0.5
VALUE_PROP_BADGE_SIZE_IN = 0.7
# Bronze for the customer side, brand green for the product side --
# matches the story template's "problem is warm/bronze, solution is green"
# convention, kept consistent across templates that pair a need with an
# answer to it.
VALUE_PROP_CUSTOMER_COLOR = "8A5A0A"
VALUE_PROP_PRODUCT_COLOR = "00754A"

POSITIONING_EYEBROW_Y_IN = 0.6
POSITIONING_CARD_TOP_IN = 1.5
POSITIONING_CARD_H_IN = 4.6
POSITIONING_CARD_MARGIN_IN = 1.0
POSITIONING_BOTTOM_BAR_H_IN = 0.14
# Cycled across the sentence's filled-in slots so each one reads as a
# distinct piece of the pitch rather than a wall of bold text.
POSITIONING_RUN_COLORS = ["00754A", "C9A227", "14B0A0", "8A5A0A"]

RACI_TITLE_Y_IN = 0.5
RACI_TABLE_TOP_IN = 1.35
RACI_TABLE_LEFT_IN = 0.7
RACI_TABLE_W_IN = 11.9
RACI_TASK_COL_W_IN = 4.2
RACI_ROLE_COL_W_IN = (RACI_TABLE_W_IN - RACI_TASK_COL_W_IN) / 4
RACI_HEADER_H_IN = 0.75
RACI_ROW_H_IN = 0.85
RACI_ROLE_LABELS = [("R", "Responsible"), ("A", "Accountable"), ("C", "Consulted"), ("I", "Informed")]
RACI_ROLE_COLORS = ["00754A", "C9A227", "14B0A0", "8A5A0A"]


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _save_bytes(prs: Presentation) -> bytes:
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def _polar(cx: float, cy: float, r: float, degrees: float) -> tuple[float, float]:
    rad = math.radians(degrees)
    return cx + r * math.cos(rad), cy + r * math.sin(rad)


def _ring_segment_points(
    cx: float, cy: float, r_in: float, r_out: float, a0: float, a1: float, steps: int = 8
) -> list[tuple[float, float]]:
    """Corner points for a thin arc (annulus sector), approximated with
    straight segments -- sidesteps autoshape angle-adjustment units
    entirely (no BLOCK_ARC-style risk), same reasoning as the pyramid's
    freeform bands."""
    points = [_polar(cx, cy, r_out, a0 + (a1 - a0) * i / steps) for i in range(steps + 1)]
    points += [_polar(cx, cy, r_in, a1 - (a1 - a0) * i / steps) for i in range(steps + 1)]
    return points


def _set_fill(shape, hex_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(hex_color)
    shape.line.fill.background()


def _add_label(slide, x_in: float, y_in: float, w_in: float, text: str, *, size: int, bold: bool, color: str, align=PP_ALIGN.CENTER, h_in: float = 0.6):
    box = slide.shapes.add_textbox(Inches(x_in - w_in / 2), Inches(y_in - h_in / 2), Inches(w_in), Inches(h_in))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = "Arial"
    return box


def _add_background(slide, prs: Presentation) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    _set_fill(bg, "FAF9F6")
    bg.shadow.inherit = False


def _add_rotated_label(slide, cx_in: float, cy_in: float, visual_h_in: float, text: str, *, size: int, color: str):
    """A label that reads bottom-to-top, for a y-axis. Built by sizing the
    box for its PRE-rotation orientation (width = the height it should
    visually occupy, height = its visual thickness) centered on the target
    point, then rotating 270 degrees -- verified round-trip via python-pptx
    that rotation is preserved and the box stays centered on (cx, cy).
    """
    w_in, h_in = visual_h_in, 0.4
    box = slide.shapes.add_textbox(Inches(cx_in - w_in / 2), Inches(cy_in - h_in / 2), Inches(w_in), Inches(h_in))
    box.rotation = 270
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.name = "Arial"


def add_wheel_slide(prs: Presentation, data: InfographicWheel) -> None:
    """Populates the fixed 5-slot radial-wheel geometry with generated content.

    Deliberately built with PIE shapes (a solid slice from the true circle
    center, verified against python-pptx's own adjustment semantics) plus a
    covering hub circle drawn last, rather than a ring/donut autoshape --
    same final look, without relying on that shape type's angle-adjustment
    units, which don't follow the same convention as PIE's and are easy to
    get subtly wrong.
    """
    # This template's geometry is fixed at exactly 5 slots; anything short
    # is padded with blank cards rather than reflowing the wheel.
    items = data.items[:5]
    items += [WheelItem(label="", description="") for _ in range(5 - len(items))]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    step = 36  # 180 degrees / 5 wedges
    start = -90
    wedge_size = Inches(2 * R_OUT_IN)
    wedge_left = Inches(WHEEL_CX_IN - R_OUT_IN)
    wedge_top = Inches(WHEEL_CY_IN - R_OUT_IN)

    for i, item in enumerate(items):
        a0 = start + i * step
        a1 = start + (i + 1) * step
        wedge = slide.shapes.add_shape(MSO_SHAPE.PIE, wedge_left, wedge_top, wedge_size, wedge_size)
        # PIE's adj1/adj2 are fractions of a full turn (verified: -0.25 -> a
        # -90 degree start), unlike some other arc-family autoshapes whose
        # angle adjustments use a different internal unit.
        wedge.adjustments[0] = a0 / 360
        wedge.adjustments[1] = a1 / 360
        _set_fill(wedge, WEDGE_COLORS[i])
        wedge.shadow.inherit = False

        mid = (a0 + a1) / 2
        lx, ly = _polar(WHEEL_CX_IN, WHEEL_CY_IN, LABEL_RADIUS_IN, mid)
        if item.label:
            _add_label(slide, lx, ly, 1.5, item.label.upper(), size=12, bold=True, color="FFFFFF")

    # Hub circle drawn last so it sits above the wedges' inner points,
    # producing the same donut-hole look as a ring-segment shape would.
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(WHEEL_CX_IN - R_IN_IN), Inches(WHEEL_CY_IN - R_IN_IN),
        Inches(2 * R_IN_IN), Inches(2 * R_IN_IN),
    )
    _set_fill(hub, "FFFFFF")
    hub.line.color.rgb = RGBColor.from_string("BFE3D3")
    hub.line.width = Pt(3)
    hub.line.fill.solid()
    hub.line.fill.fore_color.rgb = RGBColor.from_string("BFE3D3")
    hub.shadow.inherit = False
    _add_label(slide, WHEEL_CX_IN, WHEEL_CY_IN, 2 * R_IN_IN - 0.3, data.title.upper(), size=15, bold=True, color="1B1F1C", h_in=1.3)

    # Right-side numbered list.
    for i, item in enumerate(items):
        y = LIST_TOP_IN + i * LIST_GAP_IN
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(LIST_X_IN), Inches(y), Inches(LIST_W_IN), Inches(LIST_H_IN))
        card.adjustments[0] = 0.08
        _set_fill(card, WEDGE_COLORS[i])
        card.shadow.inherit = False

        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(LIST_X_IN + LIST_W_IN - 0.06), Inches(y), Inches(0.06), Inches(LIST_H_IN))
        _set_fill(accent, LIST_ACCENT_COLORS[i])
        accent.shadow.inherit = False

        num_box = slide.shapes.add_textbox(Inches(LIST_X_IN + 0.12), Inches(y), Inches(0.85), Inches(LIST_H_IN))
        tf = num_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"{i + 1:02d}"
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        run.font.name = "Arial"

        text_box = slide.shapes.add_textbox(Inches(LIST_X_IN + 1.0), Inches(y + 0.08), Inches(LIST_W_IN - 1.15), Inches(LIST_H_IN - 0.14))
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p_title = tf.paragraphs[0]
        run = p_title.add_run()
        run.text = item.label.upper()
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        run.font.name = "Arial"

        p_desc = tf.add_paragraph()
        p_desc.space_before = Pt(2)
        run = p_desc.add_run()
        run.text = item.description
        run.font.size = Pt(10.5)
        run.font.color.rgb = RGBColor.from_string("FFF7EC")
        run.font.name = "Arial"


def build_wheel_pptx(data: InfographicWheel) -> bytes:
    prs = _new_presentation()
    add_wheel_slide(prs, data)
    return _save_bytes(prs)


def add_roadmap_slide(prs: Presentation, data: InfographicRoadmap) -> None:
    """A Now/Next/Later roadmap: 3 fixed columns rendered with decreasing
    visual weight (dark to light, same hue) so closer-term horizons read
    as visually "stronger", plus a connecting rail with a dot per column
    to reinforce left-to-right progression -- the roadmap's visual
    distinction from the otherwise-similar Comparison Columns template.
    """
    columns = data.columns[:ROADMAP_COLUMN_COUNT]
    fixed_headings = ["Now", "Next", "Later"]
    columns += [
        RoadmapColumn(heading=fixed_headings[len(columns) + i], items=[])
        for i in range(ROADMAP_COLUMN_COUNT - len(columns))
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, ROADMAP_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    count = ROADMAP_COLUMN_COUNT
    available_w = SLIDE_W_IN - 2 * ROADMAP_MARGIN_IN - ROADMAP_COL_GAP_IN * (count - 1)
    col_w = available_w / count
    col_top = ROADMAP_COL_TOP_IN
    col_h = SLIDE_H_IN - col_top - 0.5
    bar_h = ROADMAP_BAR_H_IN
    col_centers = [ROADMAP_MARGIN_IN + i * (col_w + ROADMAP_COL_GAP_IN) + col_w / 2 for i in range(count)]

    rail = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(col_centers[0]), Inches(ROADMAP_RAIL_Y_IN - 0.02),
        Inches(col_centers[-1] - col_centers[0]), Inches(0.04),
    )
    _set_fill(rail, "C9C4B6")
    rail.shadow.inherit = False

    for i, cx in enumerate(col_centers):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(cx - 0.08), Inches(ROADMAP_RAIL_Y_IN - 0.08), Inches(0.16), Inches(0.16)
        )
        _set_fill(dot, ROADMAP_COLORS[i])
        dot.shadow.inherit = False

    for i, col in enumerate(columns):
        color = ROADMAP_COLORS[i]
        text_color = ROADMAP_TEXT_COLORS[i]
        col_left = ROADMAP_MARGIN_IN + i * (col_w + ROADMAP_COL_GAP_IN)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_left), Inches(col_top), Inches(col_w), Inches(col_h)
        )
        card.adjustments[0] = 0.06
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string("E4E1D8")
        card.line.width = Pt(1)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
        card.shadow.inherit = False

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(col_left), Inches(col_top), Inches(col_w), Inches(bar_h)
        )
        _set_fill(bar, color)
        bar.shadow.inherit = False
        _add_label(
            slide, col_left + col_w / 2, col_top + bar_h / 2, col_w - 0.3, col.heading.upper(),
            size=16, bold=True, color=text_color, h_in=bar_h - 0.1,
        )

        items = col.items[:ROADMAP_ITEM_COUNT]
        body_top = col_top + bar_h + 0.22
        body_h = col_h - bar_h - 0.4
        row_h = body_h / ROADMAP_ITEM_COUNT

        for j, item in enumerate(items):
            row_y = body_top + j * row_h
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(col_left + 0.28), Inches(row_y + row_h / 2 - 0.05),
                Inches(0.1), Inches(0.1),
            )
            _set_fill(dot, color)
            dot.shadow.inherit = False

            text_box = slide.shapes.add_textbox(
                Inches(col_left + 0.5), Inches(row_y), Inches(col_w - 0.75), Inches(row_h)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor.from_string("2A2E29")
            run.font.name = "Arial"


def build_roadmap_pptx(data: InfographicRoadmap) -> bytes:
    prs = _new_presentation()
    add_roadmap_slide(prs, data)
    return _save_bytes(prs)


def add_pyramid_slide(prs: Presentation, data: InfographicPyramid) -> None:
    """A vision/strategy pyramid.

    The vision statement runs as a full-width headline above the graphic
    rather than being squeezed into the pyramid's narrow apex, where a
    full sentence wouldn't reliably fit without overflowing into the
    band below it. The pyramid itself is built from freeform trapezoid
    bands with literal corner coordinates computed directly (similar
    triangles), sidestepping autoshape angle-adjustment units entirely
    rather than risking another BLOCK_ARC-style unit mismatch.
    """
    pillars = data.pillars[:PYRAMID_MAX_PILLARS]
    if len(pillars) < PYRAMID_MIN_PILLARS:
        pillars += [
            PyramidPillar(label="", description="")
            for _ in range(PYRAMID_MIN_PILLARS - len(pillars))
        ]
    total_bands = 1 + len(pillars)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, PYRAMID_HEADLINE_Y_IN, SLIDE_W_IN - 2.4, data.vision,
        size=24, bold=True, color="1B1F1C", h_in=PYRAMID_HEADLINE_H_IN,
    )

    band_h = (PYRAMID_BOTTOM_Y_IN - PYRAMID_TOP_Y_IN) / total_bands

    def width_at(j: int) -> float:
        return PYRAMID_APEX_W_IN + (PYRAMID_BASE_W_IN - PYRAMID_APEX_W_IN) * (j / total_bands)

    for i in range(total_bands):
        w0, w1 = width_at(i), width_at(i + 1)
        y0 = PYRAMID_TOP_Y_IN + i * band_h + (PYRAMID_GAP_IN / 2 if i > 0 else 0)
        y1 = PYRAMID_TOP_Y_IN + (i + 1) * band_h - (PYRAMID_GAP_IN / 2 if i < total_bands - 1 else 0)

        fb = slide.shapes.build_freeform(Inches(PYRAMID_CX_IN - w0 / 2), Inches(y0), scale=1)
        fb.add_line_segments(
            [
                (Inches(PYRAMID_CX_IN + w0 / 2), Inches(y0)),
                (Inches(PYRAMID_CX_IN + w1 / 2), Inches(y1)),
                (Inches(PYRAMID_CX_IN - w1 / 2), Inches(y1)),
            ],
            close=True,
        )
        band = fb.convert_to_shape()
        color = "1F2622" if i == 0 else COLUMN_COLORS[(i - 1) % len(COLUMN_COLORS)]
        _set_fill(band, color)
        band.shadow.inherit = False

        text_w = (w0 + w1) / 2 * 0.82
        cy = (y0 + y1) / 2

        if i == 0:
            _add_label(
                slide, PYRAMID_CX_IN, cy, text_w, "VISION",
                size=13, bold=True, color="FFFFFF", h_in=band_h - 0.15,
            )
            continue

        pillar = pillars[i - 1]
        box = slide.shapes.add_textbox(
            Inches(PYRAMID_CX_IN - text_w / 2), Inches(cy - (band_h - 0.2) / 2),
            Inches(text_w), Inches(band_h - 0.2),
        )
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p_label = tf.paragraphs[0]
        p_label.alignment = PP_ALIGN.CENTER
        run = p_label.add_run()
        run.text = pillar.label.upper()
        run.font.size = Pt(15)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        run.font.name = "Arial"

        if pillar.description:
            p_desc = tf.add_paragraph()
            p_desc.alignment = PP_ALIGN.CENTER
            p_desc.space_before = Pt(2)
            run = p_desc.add_run()
            run.text = pillar.description
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor.from_string("F2EFE8")
            run.font.name = "Arial"


def build_pyramid_pptx(data: InfographicPyramid) -> bytes:
    prs = _new_presentation()
    add_pyramid_slide(prs, data)
    return _save_bytes(prs)


def add_timeline_slide(prs: Presentation, data: InfographicTimeline) -> None:
    """4-6 dated milestones along a single horizontal line, with callout
    cards alternating above/below the line so neighboring cards don't
    overlap each other.
    """
    milestones = data.milestones[:TIMELINE_MAX_MILESTONES]
    n = max(len(milestones), 1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, TIMELINE_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(TIMELINE_LEFT_IN), Inches(TIMELINE_LINE_Y_IN - 0.02),
        Inches(TIMELINE_RIGHT_IN - TIMELINE_LEFT_IN), Inches(0.04),
    )
    _set_fill(line, "D8D3C4")
    line.shadow.inherit = False

    xs = (
        [TIMELINE_LEFT_IN + i * (TIMELINE_RIGHT_IN - TIMELINE_LEFT_IN) / (n - 1) for i in range(n)]
        if n > 1
        else [(TIMELINE_LEFT_IN + TIMELINE_RIGHT_IN) / 2]
    )

    for i, (milestone, x) in enumerate(zip(milestones, xs)):
        color = TIMELINE_COLORS[i % len(TIMELINE_COLORS)]
        above = i % 2 == 0

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x - 0.1), Inches(TIMELINE_LINE_Y_IN - 0.1), Inches(0.2), Inches(0.2)
        )
        _set_fill(dot, color)
        dot.shadow.inherit = False

        connector_top = TIMELINE_LINE_Y_IN - TIMELINE_CONNECTOR_H_IN if above else TIMELINE_LINE_Y_IN
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x - 0.01), Inches(connector_top), Inches(0.02), Inches(TIMELINE_CONNECTOR_H_IN)
        )
        _set_fill(connector, color)
        connector.shadow.inherit = False

        card_left = x - TIMELINE_CARD_W_IN / 2
        card_left = max(0.2, min(SLIDE_W_IN - TIMELINE_CARD_W_IN - 0.2, card_left))
        card_top = (
            TIMELINE_LINE_Y_IN - TIMELINE_CONNECTOR_H_IN - TIMELINE_CARD_H_IN
            if above
            else TIMELINE_LINE_Y_IN + TIMELINE_CONNECTOR_H_IN
        )

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_left), Inches(card_top), Inches(TIMELINE_CARD_W_IN), Inches(TIMELINE_CARD_H_IN)
        )
        card.adjustments[0] = 0.08
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string(color)
        card.line.width = Pt(1.5)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string(color)
        card.shadow.inherit = False

        text_box = slide.shapes.add_textbox(
            Inches(card_left + 0.15), Inches(card_top + 0.12), Inches(TIMELINE_CARD_W_IN - 0.3), Inches(TIMELINE_CARD_H_IN - 0.24)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0

        p_period = tf.paragraphs[0]
        p_period.alignment = PP_ALIGN.CENTER
        run = p_period.add_run()
        run.text = milestone.period.upper()
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.name = "Arial"

        p_label = tf.add_paragraph()
        p_label.alignment = PP_ALIGN.CENTER
        p_label.space_before = Pt(2)
        run = p_label.add_run()
        run.text = milestone.label
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("1B1F1C")
        run.font.name = "Arial"

        if milestone.description:
            p_desc = tf.add_paragraph()
            p_desc.alignment = PP_ALIGN.CENTER
            p_desc.space_before = Pt(3)
            run = p_desc.add_run()
            run.text = milestone.description
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor.from_string("5B5A52")
            run.font.name = "Arial"


def build_timeline_pptx(data: InfographicTimeline) -> bytes:
    prs = _new_presentation()
    add_timeline_slide(prs, data)
    return _save_bytes(prs)


def add_comparison_slide(prs: Presentation, data: InfographicComparison) -> None:
    """Populates a side-by-side comparison layout with 2-4 columns.

    Unlike the wheel's fixed 5-slot geometry, column count here is
    variable (driven by how many real options the source material
    describes), so column width/position is computed from the actual
    count rather than baked into fixed constants.
    """
    columns = data.columns[:COMPARISON_MAX_COLUMNS]
    count = max(len(columns), 1)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, COMPARISON_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.8,
    )

    available_w = SLIDE_W_IN - 2 * COMPARISON_MARGIN_IN - COMPARISON_COL_GAP_IN * (count - 1)
    col_w = available_w / count
    col_top = COMPARISON_COL_TOP_IN
    col_h = SLIDE_H_IN - col_top - 0.5
    bar_h = COMPARISON_BAR_H_IN

    for i, col in enumerate(columns):
        color = COLUMN_COLORS[i % len(COLUMN_COLORS)]
        col_left = COMPARISON_MARGIN_IN + i * (col_w + COMPARISON_COL_GAP_IN)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_left), Inches(col_top), Inches(col_w), Inches(col_h)
        )
        card.adjustments[0] = 0.06
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string("E4E1D8")
        card.line.width = Pt(1)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
        card.shadow.inherit = False

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(col_left), Inches(col_top), Inches(col_w), Inches(bar_h)
        )
        _set_fill(bar, color)
        bar.shadow.inherit = False
        _add_label(
            slide, col_left + col_w / 2, col_top + bar_h / 2, col_w - 0.3, col.heading.upper(),
            size=16, bold=True, color="FFFFFF", h_in=bar_h - 0.1,
        )

        points = col.points[:COMPARISON_POINT_COUNT]
        body_top = col_top + bar_h + 0.25
        body_h = col_h - bar_h - 0.45
        row_h = body_h / COMPARISON_POINT_COUNT

        for j, point in enumerate(points):
            row_y = body_top + j * row_h
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(col_left + 0.28), Inches(row_y + row_h / 2 - 0.05),
                Inches(0.1), Inches(0.1),
            )
            _set_fill(dot, color)
            dot.shadow.inherit = False

            text_box = slide.shapes.add_textbox(
                Inches(col_left + 0.5), Inches(row_y), Inches(col_w - 0.75), Inches(row_h)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = point
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor.from_string("2A2E29")
            run.font.name = "Arial"


def build_comparison_pptx(data: InfographicComparison) -> bytes:
    prs = _new_presentation()
    add_comparison_slide(prs, data)
    return _save_bytes(prs)


def add_bullet_slide(prs: Presentation, data: BulletSummarySlide) -> None:
    """The flexible fallback slide: a plain title + bullet list, for content
    that doesn't fit any of the shaped templates. Bullets are hand-drawn as
    a colored dot + text row rather than PowerPoint's native bullet-list XML,
    matching the dot-bullet style already used by the comparison/roadmap
    list rows for visual consistency.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, BULLET_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    bullets = data.bullets[:BULLET_MAX_COUNT]
    for i, bullet in enumerate(bullets):
        row_y = BULLET_LIST_TOP_IN + i * BULLET_ROW_H_IN
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(BULLET_LIST_LEFT_IN), Inches(row_y + BULLET_ROW_H_IN / 2 - 0.06),
            Inches(0.12), Inches(0.12),
        )
        _set_fill(dot, BULLET_ACCENT_COLOR)
        dot.shadow.inherit = False

        text_box = slide.shapes.add_textbox(
            Inches(BULLET_LIST_LEFT_IN + 0.4), Inches(row_y),
            Inches(BULLET_LIST_WIDTH_IN - 0.4), Inches(BULLET_ROW_H_IN),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = bullet
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor.from_string("2A2E29")
        run.font.name = "Arial"


def build_bullets_pptx(data: BulletSummarySlide) -> bytes:
    prs = _new_presentation()
    add_bullet_slide(prs, data)
    return _save_bytes(prs)


def add_matrix_slide(prs: Presentation, data: InfographicMatrix) -> None:
    """A 2x2 grid -- prioritization matrix or SWOT, depending on the axis
    labels/quadrant content the LLM chose. Axis labels use `_add_rotated_label`
    for the y-axis rather than a ring/BLOCK_ARC-style shape trick, since this
    is just plain rotated text, not a curved shape.
    """
    quadrants = data.quadrants[:MATRIX_QUADRANT_COUNT]
    quadrants += [
        MatrixQuadrant(label="", items=[])
        for _ in range(MATRIX_QUADRANT_COUNT - len(quadrants))
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, MATRIX_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    cell_w = (MATRIX_GRID_W_IN - MATRIX_GRID_GAP_IN) / 2
    cell_h = (MATRIX_GRID_H_IN - MATRIX_GRID_GAP_IN) / 2
    bar_h = MATRIX_BAR_H_IN

    # Quadrant order matches the schema: top-left, top-right, bottom-left, bottom-right.
    positions = [
        (MATRIX_GRID_LEFT_IN, MATRIX_GRID_TOP_IN),
        (MATRIX_GRID_LEFT_IN + cell_w + MATRIX_GRID_GAP_IN, MATRIX_GRID_TOP_IN),
        (MATRIX_GRID_LEFT_IN, MATRIX_GRID_TOP_IN + cell_h + MATRIX_GRID_GAP_IN),
        (MATRIX_GRID_LEFT_IN + cell_w + MATRIX_GRID_GAP_IN, MATRIX_GRID_TOP_IN + cell_h + MATRIX_GRID_GAP_IN),
    ]

    for i, (quad, (qx, qy)) in enumerate(zip(quadrants, positions)):
        color = COLUMN_COLORS[i % len(COLUMN_COLORS)]

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(qx), Inches(qy), Inches(cell_w), Inches(cell_h)
        )
        card.adjustments[0] = 0.04
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string("E4E1D8")
        card.line.width = Pt(1)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
        card.shadow.inherit = False

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(qx), Inches(qy), Inches(cell_w), Inches(bar_h)
        )
        _set_fill(bar, color)
        bar.shadow.inherit = False
        _add_label(
            slide, qx + cell_w / 2, qy + bar_h / 2, cell_w - 0.3, quad.label.upper(),
            size=14, bold=True, color="FFFFFF", h_in=bar_h - 0.1,
        )

        items = quad.items[:MATRIX_ITEM_COUNT]
        body_top = qy + bar_h + 0.15
        body_h = cell_h - bar_h - 0.3
        row_h = body_h / MATRIX_ITEM_COUNT

        for j, item in enumerate(items):
            row_y = body_top + j * row_h
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(qx + 0.2), Inches(row_y + row_h / 2 - 0.04),
                Inches(0.08), Inches(0.08),
            )
            _set_fill(dot, color)
            dot.shadow.inherit = False

            text_box = slide.shapes.add_textbox(
                Inches(qx + 0.38), Inches(row_y), Inches(cell_w - 0.55), Inches(row_h)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor.from_string("2A2E29")
            run.font.name = "Arial"

    _add_rotated_label(
        slide, MATRIX_GRID_LEFT_IN - 0.55, MATRIX_GRID_TOP_IN + MATRIX_GRID_H_IN / 2, MATRIX_GRID_H_IN * 0.6,
        data.y_axis_label, size=13, color="5B5A52",
    )
    _add_label(
        slide, MATRIX_GRID_LEFT_IN + MATRIX_GRID_W_IN / 2, MATRIX_GRID_TOP_IN + MATRIX_GRID_H_IN + 0.3,
        MATRIX_GRID_W_IN * 0.6, data.x_axis_label.upper(),
        size=13, bold=True, color="5B5A52", h_in=0.4,
    )


def build_matrix_pptx(data: InfographicMatrix) -> bytes:
    prs = _new_presentation()
    add_matrix_slide(prs, data)
    return _save_bytes(prs)


def add_story_slide(prs: Presentation, data: FeatureStory) -> None:
    """A single feature/epic's Problem -> Solution -> Impact narrative: a
    full-width headline (the one-sentence business value, kept outside the
    panels for the same reason the pyramid's vision statement is -- a real
    sentence needs more width than a narrow panel gives it) above 3 panels
    connected by arrow glyphs, so it reads as causality rather than a
    parallel comparison like comparison_columns.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, STORY_HEADLINE_Y_IN, SLIDE_W_IN - 2.0, data.headline,
        size=22, bold=True, color="1B1F1C", h_in=STORY_HEADLINE_H_IN,
    )

    panel_w = (SLIDE_W_IN - 2 * STORY_MARGIN_IN - 2 * STORY_GAP_IN) / 3
    acts = [("problem", data.problem), ("solution", data.solution), ("impact", data.impact)]

    for i, (key, act) in enumerate(acts):
        color = STORY_COLORS[key]
        px = STORY_MARGIN_IN + i * (panel_w + STORY_GAP_IN)
        py = STORY_PANEL_TOP_IN

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(px), Inches(py), Inches(panel_w), Inches(STORY_PANEL_H_IN)
        )
        card.adjustments[0] = 0.05
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string("E4E1D8")
        card.line.width = Pt(1)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
        card.shadow.inherit = False

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(px), Inches(py), Inches(panel_w), Inches(STORY_BAR_H_IN)
        )
        _set_fill(bar, color)
        bar.shadow.inherit = False
        _add_label(
            slide, px + panel_w / 2, py + STORY_BAR_H_IN / 2, panel_w - 0.3, act.heading.upper(),
            size=13, bold=True, color="FFFFFF", h_in=STORY_BAR_H_IN - 0.1,
        )

        body_box = slide.shapes.add_textbox(
            Inches(px + 0.3), Inches(py + STORY_BAR_H_IN + 0.25), Inches(panel_w - 0.6), Inches(2.1)
        )
        tf = body_box.text_frame
        tf.word_wrap = True
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = act.body
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("1B1F1C")
        run.font.name = "Arial"

        if act.detail:
            detail_box = slide.shapes.add_textbox(
                Inches(px + 0.3), Inches(py + STORY_PANEL_H_IN - 1.1), Inches(panel_w - 0.6), Inches(0.9)
            )
            tf2 = detail_box.text_frame
            tf2.word_wrap = True
            tf2.vertical_anchor = MSO_ANCHOR.BOTTOM
            tf2.margin_left = 0
            tf2.margin_right = 0
            tf2.margin_top = 0
            tf2.margin_bottom = 0
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = act.detail
            run2.font.size = Pt(11)
            run2.font.italic = True
            run2.font.color.rgb = RGBColor.from_string(color)
            run2.font.name = "Arial"

    # Arrow connectors between panels -- what makes this read as causality
    # rather than 3 parallel columns.
    for i in range(2):
        gx = STORY_MARGIN_IN + (i + 1) * panel_w + i * STORY_GAP_IN + STORY_GAP_IN / 2
        gy = STORY_PANEL_TOP_IN + STORY_PANEL_H_IN / 2
        _add_label(slide, gx, gy, STORY_GAP_IN + 0.3, "→", size=28, bold=True, color="9B968A", h_in=0.5)


def build_story_pptx(data: FeatureStory) -> bytes:
    prs = _new_presentation()
    add_story_slide(prs, data)
    return _save_bytes(prs)


def add_hub_spoke_slide(prs: Presentation, data: InfographicHubSpoke) -> None:
    """A 6-item hub-and-spoke, adapted from a stock reference design: a
    decorative ring of 6 thin color-coded arc segments around a central
    hub, each linked by a dashed connector + dot to a two-column card
    layout (3 left, 3 right) -- visually distinct from radial_wheel's
    solid pie wedges and single-side list, and closer to how the reference
    used a lighter connector language to stay legible at 6 items.
    """
    items = data.items[:HUB_SPOKE_ITEM_COUNT]
    items += [HubSpokeItem(label="", description="") for _ in range(HUB_SPOKE_ITEM_COUNT - len(items))]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    cx, cy = HUB_SPOKE_CX_IN, HUB_SPOKE_CY_IN

    total_h = 3 * HUB_SPOKE_CARD_H_IN + 2 * HUB_SPOKE_CARD_ROW_GAP_IN
    top = cy - total_h / 2
    row_cy = [top + i * (HUB_SPOKE_CARD_H_IN + HUB_SPOKE_CARD_ROW_GAP_IN) + HUB_SPOKE_CARD_H_IN / 2 for i in range(3)]

    right_x = cx + HUB_SPOKE_RING_R_OUT_IN + HUB_SPOKE_SPOKE_GAP_IN
    left_x = cx - HUB_SPOKE_RING_R_OUT_IN - HUB_SPOKE_SPOKE_GAP_IN - HUB_SPOKE_CARD_W_IN

    # items[0:3] = left column (top to bottom), items[3:6] = right column.
    slots = [(i, "left", left_x) for i in range(3)] + [(i, "right", right_x) for i in range(3)]

    for idx, (item, (row_i, side, card_x)) in enumerate(zip(items, slots)):
        color = HUB_SPOKE_COLORS[idx % len(HUB_SPOKE_COLORS)]
        card_y = row_cy[row_i] - HUB_SPOKE_CARD_H_IN / 2

        near_x = card_x + HUB_SPOKE_CARD_W_IN if side == "left" else card_x
        near_y = row_cy[row_i]
        angle = math.degrees(math.atan2(near_y - cy, near_x - cx))

        pts = _ring_segment_points(
            cx, cy, HUB_SPOKE_RING_R_IN_IN, HUB_SPOKE_RING_R_OUT_IN,
            angle - HUB_SPOKE_RING_HALF_SPAN_DEG, angle + HUB_SPOKE_RING_HALF_SPAN_DEG,
        )
        fb = slide.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]), scale=1)
        fb.add_line_segments([(Inches(x), Inches(y)) for x, y in pts[1:]], close=True)
        ring_shape = fb.convert_to_shape()
        _set_fill(ring_shape, color)
        ring_shape.shadow.inherit = False

        edge_x, edge_y = _polar(cx, cy, HUB_SPOKE_RING_R_OUT_IN, angle)
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(edge_x), Inches(edge_y), Inches(near_x), Inches(near_y)
        )
        connector.line.color.rgb = RGBColor.from_string(color)
        connector.line.width = Pt(1.5)
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH

        dot_x, dot_y = (edge_x + near_x) / 2, (edge_y + near_y) / 2
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(dot_x - 0.06), Inches(dot_y - 0.06), Inches(0.12), Inches(0.12)
        )
        _set_fill(dot, color)
        dot.shadow.inherit = False

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(card_y), Inches(HUB_SPOKE_CARD_W_IN), Inches(HUB_SPOKE_CARD_H_IN)
        )
        card.adjustments[0] = 0.5
        _set_fill(card, "FFFFFF")
        card.line.color.rgb = RGBColor.from_string("E4E1D8")
        card.line.width = Pt(1)
        card.line.fill.solid()
        card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
        card.shadow.inherit = False

        cap_x = card_x if side == "right" else card_x + HUB_SPOKE_CARD_W_IN - HUB_SPOKE_CAP_W_IN
        cap = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(cap_x), Inches(card_y), Inches(HUB_SPOKE_CAP_W_IN), Inches(HUB_SPOKE_CARD_H_IN)
        )
        _set_fill(cap, color)
        cap.shadow.inherit = False
        _add_label(
            slide, cap_x + HUB_SPOKE_CAP_W_IN / 2, card_y + HUB_SPOKE_CARD_H_IN / 2, HUB_SPOKE_CAP_W_IN - 0.2,
            f"{idx + 1:02d}", size=22, bold=True, color="FFFFFF", h_in=HUB_SPOKE_CARD_H_IN - 0.2,
        )

        text_x = card_x + HUB_SPOKE_CAP_W_IN + 0.2 if side == "right" else card_x + 0.2
        text_w = HUB_SPOKE_CARD_W_IN - HUB_SPOKE_CAP_W_IN - 0.4
        text_box = slide.shapes.add_textbox(
            Inches(text_x), Inches(card_y + 0.12), Inches(text_w), Inches(HUB_SPOKE_CARD_H_IN - 0.24)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p_label = tf.paragraphs[0]
        run = p_label.add_run()
        run.text = item.label.upper()
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("1B1F1C")
        run.font.name = "Arial"

        if item.description:
            p_desc = tf.add_paragraph()
            p_desc.space_before = Pt(2)
            run = p_desc.add_run()
            run.text = item.description
            run.font.size = Pt(10)
            run.font.color.rgb = RGBColor.from_string("5B5A52")
            run.font.name = "Arial"

    # Hub drawn last so it sits above the ring segments' inner points.
    hub = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - HUB_SPOKE_HUB_R_IN), Inches(cy - HUB_SPOKE_HUB_R_IN),
        Inches(2 * HUB_SPOKE_HUB_R_IN), Inches(2 * HUB_SPOKE_HUB_R_IN),
    )
    _set_fill(hub, "FFFFFF")
    hub.line.color.rgb = RGBColor.from_string("E4E1D8")
    hub.line.width = Pt(2)
    hub.line.fill.solid()
    hub.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
    hub.shadow.inherit = False

    hub_text = slide.shapes.add_textbox(
        Inches(cx - HUB_SPOKE_HUB_R_IN + 0.15), Inches(cy - HUB_SPOKE_HUB_R_IN + 0.15),
        Inches(2 * HUB_SPOKE_HUB_R_IN - 0.3), Inches(2 * HUB_SPOKE_HUB_R_IN - 0.3),
    )
    tf = hub_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.CENTER
    run = p_title.add_run()
    run.text = data.title.upper()
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("1B1F1C")
    run.font.name = "Arial"

    if data.description:
        p_desc = tf.add_paragraph()
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.space_before = Pt(4)
        run = p_desc.add_run()
        run.text = data.description
        run.font.size = Pt(9.5)
        run.font.color.rgb = RGBColor.from_string("5B5A52")
        run.font.name = "Arial"


def build_hub_spoke_pptx(data: InfographicHubSpoke) -> bytes:
    prs = _new_presentation()
    add_hub_spoke_slide(prs, data)
    return _save_bytes(prs)


def add_title_slide(prs: Presentation, data: TitleSlide) -> None:
    """The deck's cover slide: a thin brand-green top bar to mark it as the
    opener (content slides carry no such bar), a centered headline + gold
    rule + subtitle, and a row of highlight chips -- distinct from
    bullet_summary's left-aligned title-and-list shape on purpose, since
    this isn't a content section.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(TITLE_TOP_BAR_H_IN)
    )
    _set_fill(top_bar, BRAND_COLORS[0])
    top_bar.shadow.inherit = False

    _add_label(
        slide, SLIDE_W_IN / 2, TITLE_HEADLINE_Y_IN, SLIDE_W_IN - 2.0, data.title,
        size=40, bold=True, color="1F2622", h_in=TITLE_HEADLINE_H_IN,
    )

    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(SLIDE_W_IN / 2 - TITLE_RULE_W_IN / 2), Inches(TITLE_RULE_Y_IN),
        Inches(TITLE_RULE_W_IN), Inches(TITLE_RULE_H_IN),
    )
    _set_fill(rule, BRAND_COLORS[1])
    rule.shadow.inherit = False

    _add_label(
        slide, SLIDE_W_IN / 2, TITLE_SUBTITLE_Y_IN, TITLE_SUBTITLE_W_IN, data.subtitle,
        size=18, bold=False, color="5B5A52", h_in=TITLE_SUBTITLE_H_IN,
    )

    highlights = data.highlights[:TITLE_HIGHLIGHT_MAX]
    count = len(highlights)
    if count:
        total_w = count * TITLE_CHIP_W_IN + (count - 1) * TITLE_CHIP_GAP_IN
        start_x = SLIDE_W_IN / 2 - total_w / 2
        for i, highlight in enumerate(highlights):
            chip_x = start_x + i * (TITLE_CHIP_W_IN + TITLE_CHIP_GAP_IN)
            color = BRAND_COLORS[i % len(BRAND_COLORS)]
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(chip_x), Inches(TITLE_CHIP_Y_IN), Inches(TITLE_CHIP_W_IN), Inches(TITLE_CHIP_H_IN),
            )
            chip.adjustments[0] = 0.5
            _set_fill(chip, color)
            chip.shadow.inherit = False
            _add_label(
                slide, chip_x + TITLE_CHIP_W_IN / 2, TITLE_CHIP_Y_IN + TITLE_CHIP_H_IN / 2,
                TITLE_CHIP_W_IN - 0.2, highlight.upper(),
                size=12, bold=True, color="FFFFFF", h_in=TITLE_CHIP_H_IN - 0.1,
            )


def build_title_pptx(data: TitleSlide) -> bytes:
    prs = _new_presentation()
    add_title_slide(prs, data)
    return _save_bytes(prs)


def add_agenda_slide(prs: Presentation, data: AgendaSlide) -> None:
    """A table-of-contents slide: one row per other slide in the deck, each
    with a colored page-number chip (the literal page it lands on, not just
    a running count) and a divider hairline -- deliberately close to
    bullet_summary's dot-and-text row, but the page number is real
    information here, so it replaces the plain dot.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, AGENDA_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    for i, item in enumerate(data.items):
        row_y = AGENDA_LIST_TOP_IN + i * AGENDA_ROW_H_IN
        color = BRAND_COLORS[i % len(BRAND_COLORS)]

        chip = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(AGENDA_LIST_LEFT_IN), Inches(row_y + (AGENDA_ROW_H_IN - AGENDA_CHIP_SIZE_IN) / 2),
            Inches(AGENDA_CHIP_SIZE_IN), Inches(AGENDA_CHIP_SIZE_IN),
        )
        chip.adjustments[0] = 0.2
        _set_fill(chip, color)
        chip.shadow.inherit = False
        _add_label(
            slide, AGENDA_LIST_LEFT_IN + AGENDA_CHIP_SIZE_IN / 2, row_y + AGENDA_ROW_H_IN / 2,
            AGENDA_CHIP_SIZE_IN - 0.05, f"{item.page:02d}",
            size=14, bold=True, color="FFFFFF", h_in=AGENDA_CHIP_SIZE_IN - 0.05,
        )

        label_x = AGENDA_LIST_LEFT_IN + AGENDA_CHIP_SIZE_IN + 0.35
        text_box = slide.shapes.add_textbox(
            Inches(label_x), Inches(row_y),
            Inches(AGENDA_LIST_WIDTH_IN - AGENDA_CHIP_SIZE_IN - 0.35), Inches(AGENDA_ROW_H_IN),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = item.label
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("2A2E29")
        run.font.name = "Arial"

        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(AGENDA_LIST_LEFT_IN), Inches(row_y + AGENDA_ROW_H_IN - 0.01),
            Inches(AGENDA_LIST_WIDTH_IN), Pt(1),
        )
        _set_fill(divider, "E4E1D8")
        divider.shadow.inherit = False


def build_agenda_pptx(data: AgendaSlide) -> bytes:
    prs = _new_presentation()
    add_agenda_slide(prs, data)
    return _save_bytes(prs)


def _draw_value_prop_panel(
    slide, x_in: float, panel_label: str, sections: list[tuple[str, list[str]]], header_color: str
) -> None:
    """One side of the value-proposition canvas: a header bar naming the
    side (Customer / Our Product), then 3 stacked labeled mini-lists.
    Shared by both sides so their spacing math can't drift apart."""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x_in), Inches(VALUE_PROP_PANEL_TOP_IN), Inches(VALUE_PROP_PANEL_W_IN), Inches(VALUE_PROP_PANEL_H_IN),
    )
    card.adjustments[0] = 0.03
    _set_fill(card, "FFFFFF")
    card.line.color.rgb = RGBColor.from_string("E4E1D8")
    card.line.width = Pt(1)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
    card.shadow.inherit = False

    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x_in), Inches(VALUE_PROP_PANEL_TOP_IN), Inches(VALUE_PROP_PANEL_W_IN), Inches(VALUE_PROP_HEADER_H_IN),
    )
    _set_fill(header, header_color)
    header.shadow.inherit = False
    _add_label(
        slide, x_in + VALUE_PROP_PANEL_W_IN / 2, VALUE_PROP_PANEL_TOP_IN + VALUE_PROP_HEADER_H_IN / 2,
        VALUE_PROP_PANEL_W_IN - 0.3, panel_label,
        size=15, bold=True, color="FFFFFF", h_in=VALUE_PROP_HEADER_H_IN - 0.1,
    )

    body_top = VALUE_PROP_PANEL_TOP_IN + VALUE_PROP_HEADER_H_IN
    body_h = VALUE_PROP_PANEL_H_IN - VALUE_PROP_HEADER_H_IN
    section_h = body_h / len(sections)
    label_h = 0.3
    text_left = x_in + 0.3
    text_w = VALUE_PROP_PANEL_W_IN - 0.55

    for i, (label, items) in enumerate(sections):
        sy = body_top + i * section_h

        label_box = slide.shapes.add_textbox(Inches(text_left), Inches(sy + 0.08), Inches(text_w), Inches(label_h))
        tf = label_box.text_frame
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor.from_string(header_color)
        run.font.name = "Arial"

        items = items[:VALUE_PROP_MAX_ITEMS]
        rows_top = sy + label_h + 0.1
        rows_h = section_h - label_h - 0.15
        row_h = rows_h / VALUE_PROP_MAX_ITEMS
        for j, item in enumerate(items):
            row_y = rows_top + j * row_h
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, Inches(text_left), Inches(row_y + row_h / 2 - 0.035), Inches(0.07), Inches(0.07),
            )
            _set_fill(dot, header_color)
            dot.shadow.inherit = False

            text_box = slide.shapes.add_textbox(
                Inches(text_left + 0.2), Inches(row_y), Inches(text_w - 0.2), Inches(row_h)
            )
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item
            run.font.size = Pt(10.5)
            run.font.color.rgb = RGBColor.from_string("2A2E29")
            run.font.name = "Arial"


def add_value_proposition_slide(prs: Presentation, data: ValuePropositionSlide) -> None:
    """A Value Proposition Canvas: customer jobs/pains/gains on the left,
    the product's offerings/relievers/creators on the right, with a small
    badge marking the fit between them -- the business-value argument made
    by showing which need each part of the product answers, rather than a
    metric asserted on its own.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, VALUE_PROP_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    _draw_value_prop_panel(
        slide, VALUE_PROP_LEFT_X_IN, "CUSTOMER",
        [("JOBS", data.customer_jobs), ("PAINS", data.customer_pains), ("GAINS", data.customer_gains)],
        VALUE_PROP_CUSTOMER_COLOR,
    )
    _draw_value_prop_panel(
        slide, VALUE_PROP_RIGHT_X_IN, "OUR PRODUCT",
        [("PRODUCTS & SERVICES", data.products_services), ("PAIN RELIEVERS", data.pain_relievers), ("GAIN CREATORS", data.gain_creators)],
        VALUE_PROP_PRODUCT_COLOR,
    )

    badge_cx = SLIDE_W_IN / 2
    badge_cy = VALUE_PROP_PANEL_TOP_IN + VALUE_PROP_PANEL_H_IN / 2
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(badge_cx - VALUE_PROP_BADGE_SIZE_IN / 2), Inches(badge_cy - VALUE_PROP_BADGE_SIZE_IN / 2),
        Inches(VALUE_PROP_BADGE_SIZE_IN), Inches(VALUE_PROP_BADGE_SIZE_IN),
    )
    _set_fill(badge, BRAND_COLORS[1])
    badge.line.color.rgb = RGBColor.from_string("FFFFFF")
    badge.line.width = Pt(3)
    badge.line.fill.solid()
    badge.line.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
    badge.shadow.inherit = False
    _add_label(
        slide, badge_cx, badge_cy, VALUE_PROP_BADGE_SIZE_IN - 0.1, "FIT",
        size=13, bold=True, color="FFFFFF", h_in=VALUE_PROP_BADGE_SIZE_IN - 0.1,
    )


def build_value_proposition_pptx(data: ValuePropositionSlide) -> bytes:
    prs = _new_presentation()
    add_value_proposition_slide(prs, data)
    return _save_bytes(prs)


def add_positioning_statement_slide(prs: Presentation, data: PositioningStatementSlide) -> None:
    """The Geoffrey Moore positioning statement, assembled as one sentence
    with the filled-in slots color-emphasized -- deliberately a single
    flowing narrative (multiple runs in one paragraph) rather than a form
    of labeled fields, since the point is that it reads aloud as a pitch.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, POSITIONING_EYEBROW_Y_IN, SLIDE_W_IN - 1.0, "POSITIONING STATEMENT",
        size=14, bold=True, color="8A8781", h_in=0.4,
    )

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(POSITIONING_CARD_MARGIN_IN), Inches(POSITIONING_CARD_TOP_IN),
        Inches(SLIDE_W_IN - 2 * POSITIONING_CARD_MARGIN_IN), Inches(POSITIONING_CARD_H_IN),
    )
    card.adjustments[0] = 0.04
    _set_fill(card, "FFFFFF")
    card.line.color.rgb = RGBColor.from_string("E4E1D8")
    card.line.width = Pt(1)
    card.line.fill.solid()
    card.line.fill.fore_color.rgb = RGBColor.from_string("E4E1D8")
    card.shadow.inherit = False

    inner_margin = 0.7
    text_box = slide.shapes.add_textbox(
        Inches(POSITIONING_CARD_MARGIN_IN + inner_margin), Inches(POSITIONING_CARD_TOP_IN + 0.4),
        Inches(SLIDE_W_IN - 2 * (POSITIONING_CARD_MARGIN_IN + inner_margin)), Inches(POSITIONING_CARD_H_IN - 0.8),
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    ink = "2A2E29"
    plain_size = Pt(20)
    bold_size = Pt(22)
    segments = [
        ("For ", False, ink),
        (data.target_customer, True, POSITIONING_RUN_COLORS[0]),
        (" who ", False, ink),
        (data.need, True, POSITIONING_RUN_COLORS[1]),
        (", ", False, ink),
        (data.product_name, True, POSITIONING_RUN_COLORS[3]),
        (" is a ", False, ink),
        (data.category, True, POSITIONING_RUN_COLORS[0]),
        (" that ", False, ink),
        (data.key_benefit, True, POSITIONING_RUN_COLORS[1]),
        (". Unlike ", False, ink),
        (data.primary_alternative, True, POSITIONING_RUN_COLORS[2]),
        (", we ", False, ink),
        (data.differentiator, True, POSITIONING_RUN_COLORS[3]),
        (".", False, ink),
    ]
    for text, bold, color in segments:
        run = p.add_run()
        run.text = text
        run.font.size = bold_size if bold else plain_size
        run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color)
        run.font.name = "Arial"

    bottom_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(SLIDE_H_IN - POSITIONING_BOTTOM_BAR_H_IN),
        prs.slide_width, Inches(POSITIONING_BOTTOM_BAR_H_IN),
    )
    _set_fill(bottom_bar, BRAND_COLORS[0])
    bottom_bar.shadow.inherit = False


def build_positioning_statement_pptx(data: PositioningStatementSlide) -> bytes:
    prs = _new_presentation()
    add_positioning_statement_slide(prs, data)
    return _save_bytes(prs)


def add_raci_chart_slide(prs: Presentation, data: RaciChartSlide) -> None:
    """Who owns what: a task column plus 4 fixed RACI role columns, each
    cell holding a name/role rather than a matrix of R/A/C/I letters --
    reads directly as "who do I talk to" without a legend.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide, prs)

    _add_label(
        slide, SLIDE_W_IN / 2, RACI_TITLE_Y_IN, SLIDE_W_IN - 1.0, data.title.upper(),
        size=26, bold=True, color="1B1F1C", h_in=0.7,
    )

    col_x = [RACI_TABLE_LEFT_IN + RACI_TASK_COL_W_IN + i * RACI_ROLE_COL_W_IN for i in range(5)]

    # Header row: plain task-column label, then 4 colored role cells.
    _add_label(
        slide, RACI_TABLE_LEFT_IN + RACI_TASK_COL_W_IN / 2, RACI_TABLE_TOP_IN + RACI_HEADER_H_IN / 2,
        RACI_TASK_COL_W_IN - 0.2, "INITIATIVE",
        size=13, bold=True, color="5B5A52", align=PP_ALIGN.LEFT, h_in=RACI_HEADER_H_IN - 0.1,
    )
    for i, (letter, word) in enumerate(RACI_ROLE_LABELS):
        color = RACI_ROLE_COLORS[i]
        cell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(col_x[i] + 0.05), Inches(RACI_TABLE_TOP_IN),
            Inches(RACI_ROLE_COL_W_IN - 0.1), Inches(RACI_HEADER_H_IN),
        )
        cell.adjustments[0] = 0.12
        _set_fill(cell, color)
        cell.shadow.inherit = False
        _add_label(
            slide, col_x[i] + RACI_ROLE_COL_W_IN / 2, RACI_TABLE_TOP_IN + 0.27, RACI_ROLE_COL_W_IN - 0.2, letter,
            size=18, bold=True, color="FFFFFF", h_in=0.32,
        )
        _add_label(
            slide, col_x[i] + RACI_ROLE_COL_W_IN / 2, RACI_TABLE_TOP_IN + 0.58, RACI_ROLE_COL_W_IN - 0.1, word,
            size=9, bold=False, color="FFFFFF", h_in=0.24,
        )

    rows = data.rows[:RACI_MAX_ROWS]
    body_top = RACI_TABLE_TOP_IN + RACI_HEADER_H_IN + 0.1

    for r, row in enumerate(rows):
        row_y = body_top + r * RACI_ROW_H_IN
        if r % 2 == 0:
            stripe = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(RACI_TABLE_LEFT_IN), Inches(row_y), Inches(RACI_TABLE_W_IN), Inches(RACI_ROW_H_IN - 0.06),
            )
            _set_fill(stripe, "F5F4F0")
            stripe.shadow.inherit = False

        _add_label(
            slide, RACI_TABLE_LEFT_IN + RACI_TASK_COL_W_IN / 2, row_y + (RACI_ROW_H_IN - 0.06) / 2,
            RACI_TASK_COL_W_IN - 0.3, row.task,
            size=13, bold=True, color="1F2622", align=PP_ALIGN.LEFT, h_in=RACI_ROW_H_IN - 0.16,
        )
        for i, value in enumerate([row.responsible, row.accountable, row.consulted, row.informed]):
            _add_label(
                slide, col_x[i] + RACI_ROLE_COL_W_IN / 2, row_y + (RACI_ROW_H_IN - 0.06) / 2,
                RACI_ROLE_COL_W_IN - 0.15, value,
                size=11.5, bold=False, color="2A2E29", h_in=RACI_ROW_H_IN - 0.16,
            )

    table_bottom = body_top + len(rows) * RACI_ROW_H_IN
    for i in range(5):
        divider_x = RACI_TABLE_LEFT_IN + RACI_TASK_COL_W_IN + i * RACI_ROLE_COL_W_IN
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(divider_x), Inches(body_top), Pt(1), Inches(table_bottom - body_top),
        )
        _set_fill(divider, "E4E1D8")
        divider.shadow.inherit = False


def build_raci_chart_pptx(data: RaciChartSlide) -> bytes:
    prs = _new_presentation()
    add_raci_chart_slide(prs, data)
    return _save_bytes(prs)


_DECK_SLIDE_BUILDERS: dict[str, Callable[[Presentation, InfographicDiagram], None]] = {
    "radial_wheel": add_wheel_slide,
    "comparison_columns": add_comparison_slide,
    "now_next_later": add_roadmap_slide,
    "vision_pyramid": add_pyramid_slide,
    "quarterly_timeline": add_timeline_slide,
    "bullet_summary": add_bullet_slide,
    "matrix_2x2": add_matrix_slide,
    "feature_story": add_story_slide,
    "hub_spoke": add_hub_spoke_slide,
    "title_intro": add_title_slide,
    "agenda": add_agenda_slide,
    "value_proposition": add_value_proposition_slide,
    "positioning_statement": add_positioning_statement_slide,
    "raci_chart": add_raci_chart_slide,
}


def build_deck_pptx(slides: list[InfographicDiagram]) -> bytes:
    """Assembles every planned slide into a single deck, in order. Each
    slide keeps its own template's exact visual design -- this only changes
    the outer scaffolding from "one Presentation per template" to "one
    Presentation shared across all of them."
    """
    prs = _new_presentation()
    for slide_data in slides:
        _DECK_SLIDE_BUILDERS[slide_data.template](prs, slide_data)
    return _save_bytes(prs)
